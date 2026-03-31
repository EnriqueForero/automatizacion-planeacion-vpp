# -*- coding: utf-8 -*-
"""
Generadores de PPTX y constructores de plantillas para las 3 familias.

Constructores de plantilla: crean los .pptx con marcadores {{...}}.
Generadores: leen datos y reemplazan marcadores en las plantillas.
"""

from __future__ import annotations

import os
import re
from typing import Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .config import (
    Config, FamiliaUnidad,
    EJES_REFERENCIA, ORDEN_EJES,
)
from .utils import _trunc, log


# ═══════════════════════════════════════════════════════════════════════
# GENERADOR PPTX — MISIONAL
# ═══════════════════════════════════════════════════════════════════════

class GeneradorPPTXMisional:
    """
    Genera presentaciones PPTX para unidades MISIONALES
    usando una plantilla con marcadores {{...}}.
    """

    def __init__(self, config: Config):
        self.cfg = config

    @staticmethod
    def _linea_tiene_contenido(le: Dict) -> bool:
        """Verifica si una línea tiene datos diligenciados."""
        if le.get('nombre'):
            return True
        if any(a.get('accion') or a.get('avance')
               for a in le.get('acciones', [])):
            return True
        if any(i.get('indicador') for i in le.get('indicadores', [])):
            return True
        return False

    def _construir_reemplazos(self, data: Dict) -> Dict[str, str]:
        """Construye dict marcador→valor. Vacíos quedan como ''."""
        cfg = data['config']
        q = cfg.get('Trimestre en seguimiento', 'Q1')
        unidad = cfg.get('Nombre de la unidad', 'Área')
        año = cfg.get('Año', '2026')
        tipo = cfg.get('Tipo de unidad', 'EJE')

        r = {
            '{{TRIMESTRE}}': q,
            '{{UNIDAD}}': unidad,
            '{{AÑO}}': año,
            '{{TIPO}}': tipo,
        }

        # DOFA
        prefix_map = {'DEBILIDADES': 'DEB', 'OPORTUNIDADES': 'OPO',
                      'FORTALEZAS': 'FOR', 'AMENAZAS': 'AME'}
        for cuadrante, prefix in prefix_map.items():
            items = data['dofa'].get(cuadrante, [])
            for i in range(1, self.cfg.max_dofa_por_cuadrante + 1):
                tag = f'{{{{{prefix}_BASE_{i}}}}}'
                if i <= len(items):
                    item = items[i - 1]
                    if item['estado'] == 'Se elimina':
                        r[tag] = ''
                    elif item['estado'] == 'Se actualiza' \
                            and item['actualizacion']:
                        r[tag] = _trunc(item['actualizacion'], 200)
                    else:
                        r[tag] = _trunc(item['base'], 200)
                else:
                    r[tag] = ''

        # Tendencias
        for i in range(1, self.cfg.max_tendencias + 1):
            tag = f'{{{{TEND_{i}}}}}'
            if i <= len(data['tendencias']):
                t = data['tendencias'][i - 1]
                texto = (t['actualizacion']
                         if t['estado'] == 'Se actualiza'
                         and t['actualizacion']
                         else t['base'])
                r[tag] = _trunc(texto, 300)
            else:
                r[tag] = ''

        # Líneas estratégicas (LE1 a LE5)
        max_le = self.cfg.max_lineas_estrategicas
        max_acc = self.cfg.max_acciones_por_linea

        for le_idx in range(1, max_le + 1):
            le_data = (data['lineas'][le_idx - 1]
                       if le_idx <= len(data['lineas']) else None)

            r[f'{{{{LE{le_idx}_NOMBRE}}}}'] = (
                _trunc(le_data['nombre'], 200) if le_data else ''
            )

            for acc_idx in range(1, max_acc + 1):
                tag_acc = f'{{{{LE{le_idx}_ACC_{acc_idx}}}}}'
                tag_act = f'{{{{LE{le_idx}_ACT_{acc_idx}}}}}'
                tag_ava = f'{{{{LE{le_idx}_AVA_{acc_idx}}}}}'
                if le_data and acc_idx <= len(le_data['acciones']):
                    acc = le_data['acciones'][acc_idx - 1]
                    r[tag_acc] = _trunc(acc['accion'], 200)
                    r[tag_act] = _trunc(acc.get('actividad', ''), 200)
                    r[tag_ava] = _trunc(acc['avance'],
                                       self.cfg.max_chars_slide)
                else:
                    r[tag_acc] = ''
                    r[tag_act] = ''
                    r[tag_ava] = ''

            for ind_idx in range(1, self.cfg.max_indicadores_por_linea + 1):
                tag_ind = f'{{{{LE{le_idx}_IND_{ind_idx}}}}}'
                tag_meta = f'{{{{LE{le_idx}_META_{ind_idx}}}}}'
                tag_res = f'{{{{LE{le_idx}_RES_{ind_idx}}}}}'
                if le_data and ind_idx <= len(le_data['indicadores']):
                    ind = le_data['indicadores'][ind_idx - 1]
                    r[tag_ind] = ind['indicador']
                    r[tag_meta] = str(ind['meta'])
                    r[tag_res] = str(ind['avance'])
                else:
                    r[tag_ind] = ''
                    r[tag_meta] = ''
                    r[tag_res] = ''

        # Casos de éxito
        for i in range(1, self.cfg.max_casos_exito + 1):
            tag_tit = f'{{{{CASO_{i}_TIT}}}}'
            tag_desc = f'{{{{CASO_{i}_DESC}}}}'
            if i <= len(data['casos_exito']):
                caso = data['casos_exito'][i - 1]
                r[tag_tit] = caso['titulo']
                r[tag_desc] = _trunc(caso['descripcion'],
                                    self.cfg.max_chars_slide)
            else:
                r[tag_tit] = ''
                r[tag_desc] = ''

        # Metas
        for i in range(1, self.cfg.max_metas + 1):
            tag_ind = f'{{{{META_{i}_IND}}}}'
            tag_meta = f'{{{{META_{i}_META}}}}'
            tag_ava = f'{{{{META_{i}_AVA}}}}'
            if i <= len(data['metas']):
                m = data['metas'][i - 1]
                r[tag_ind] = m['indicador']
                r[tag_meta] = str(m['meta'])
                r[tag_ava] = str(m['avance'])
            else:
                r[tag_ind] = ''
                r[tag_meta] = ''
                r[tag_ava] = ''

        return r

    @staticmethod
    def _iter_shapes(shapes):
        """Itera shapes recursivamente incluyendo grupos."""
        for shape in shapes:
            yield shape
            if shape.shape_type == 6:  # GROUP
                yield from GeneradorPPTXMisional._iter_shapes(shape.shapes)

    @staticmethod
    def _paragraph_full_text(paragraph) -> str:
        return ''.join(r.text for r in paragraph.runs)

    def _slide_tiene_marcadores_le(self, slide, le_num: int) -> bool:
        """Verifica si un slide tiene marcadores de una LE específica."""
        patron = f'LE{le_num}_'
        for shape in self._iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if patron in self._paragraph_full_text(paragraph):
                    return True
        return False

    def _reemplazar_en_slide(self, slide, reemplazos: Dict) -> int:
        """
        Reemplaza marcadores en un slide.
        Estrategia robusta: concatena runs → reemplaza → redistribuye.
        Limpia marcadores residuales y artefactos de formato vacío.
        """
        total = 0
        for shape in self._iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                runs = paragraph.runs
                if not runs:
                    continue

                full_text = ''.join(r.text for r in runs)
                if '{{' not in full_text:
                    continue

                new_text = full_text
                for marcador, valor in reemplazos.items():
                    if marcador in new_text:
                        new_text = new_text.replace(marcador, valor)
                        total += 1

                # Limpiar cualquier marcador residual → vacío
                new_text = re.sub(r'\{\{[^}]+\}\}', '', new_text)

                # Limpiar artefactos de formato vacío:
                #  " () " → "" (paréntesis vacíos de indicadores sin datos)
                new_text = re.sub(r'\s*\(\s*\)\s*', '', new_text)
                # Normalizar espacios múltiples y strip
                new_text = re.sub(r'  +', ' ', new_text).strip()

                if new_text != full_text:
                    runs[0].text = new_text
                    for r in runs[1:]:
                        r.text = ''
        return total

    @staticmethod
    def _eliminar_slides(prs, indices: List[int]):
        """Elimina slides por índice manipulando XML."""
        if not indices:
            return
        R_NS = ('{http://schemas.openxmlformats.org/'
                'officeDocument/2006/relationships}id')
        sldIdLst = prs.slides._sldIdLst
        for idx in sorted(indices, reverse=True):
            elems = list(sldIdLst)
            if idx < len(elems):
                elem = elems[idx]
                rId = elem.get(R_NS)
                sldIdLst.remove(elem)
                if rId:
                    try:
                        prs.part.drop_rel(rId)
                    except Exception as e:
                        log.warning(f"  ⚠️ Relación {rId}: {e}")
                log.info(f"  🗑️  Slide {idx + 1} eliminado")

    def generar(self, data: Dict, ruta_plantilla: str,
                ruta_salida: str) -> str:
        """Genera PPTX reemplazando marcadores en la plantilla."""
        log.info(f"📊 Plantilla: {os.path.basename(ruta_plantilla)}")
        prs = Presentation(ruta_plantilla)
        reemplazos = self._construir_reemplazos(data)
        log.info(f"  Marcadores definidos: {len(reemplazos)}")

        # Determinar líneas con contenido
        lineas_con_contenido = set()
        for le_idx in range(1, self.cfg.max_lineas_estrategicas + 1):
            if le_idx <= len(data['lineas']):
                if self._linea_tiene_contenido(data['lineas'][le_idx - 1]):
                    lineas_con_contenido.add(le_idx)
        log.info(f"  Líneas activas: {sorted(lineas_con_contenido)}")

        # Identificar slides a eliminar
        slides_a_eliminar = []
        for slide_idx, slide in enumerate(prs.slides):
            le_en_slide = set()
            for le_num in range(1, self.cfg.max_lineas_estrategicas + 1):
                if self._slide_tiene_marcadores_le(slide, le_num):
                    le_en_slide.add(le_num)
            if le_en_slide and le_en_slide.isdisjoint(lineas_con_contenido):
                slides_a_eliminar.append(slide_idx)

        # Reemplazar marcadores
        total = sum(
            self._reemplazar_en_slide(sl, reemplazos)
            for sl in prs.slides
        )
        log.info(f"  Marcadores reemplazados: {total}")

        # Eliminar slides vacíos
        if slides_a_eliminar:
            log.info(f"  Eliminando {len(slides_a_eliminar)} slides: "
                     f"{[s + 1 for s in slides_a_eliminar]}")
            self._eliminar_slides(prs, slides_a_eliminar)

        prs.save(ruta_salida)
        log.info(f"  ✅ PPTX: {os.path.basename(ruta_salida)} "
                 f"({len(prs.slides)} slides)")
        return ruta_salida


# ═══════════════════════════════════════════════════════════════════════
# CONSTRUCTOR PLANTILLA PPTX — TERRITORIAL
# ═══════════════════════════════════════════════════════════════════════

class ConstructorPlantillaTerritorial:
    """
    Construye programáticamente la plantilla PPTX territorial
    con 28 slides y todos los marcadores {{...}} necesarios.
    """

    # Colores institucionales
    C_AZ  = RGBColor(0x1B, 0x3A, 0x5C)
    C_AZ2 = RGBColor(0x2E, 0x6D, 0xA4)
    C_RJ  = RGBColor(0xC0, 0x39, 0x2B)
    C_VD  = RGBColor(0x27, 0xAE, 0x60)
    C_NJ  = RGBColor(0xF3, 0x9C, 0x12)
    C_GR  = RGBColor(0x58, 0x58, 0x58)
    C_BL  = RGBColor(0xFF, 0xFF, 0xFF)
    C_GC  = RGBColor(0x99, 0x99, 0x99)

    EJE_COLORS = {
        'MP':  RGBColor(0xE6, 0x7E, 0x22),
        'TUR': RGBColor(0x29, 0x80, 0xB9),
        'INV': RGBColor(0x27, 0xAE, 0x60),
        'EXP': RGBColor(0xC0, 0x39, 0x2B),
    }

    def __init__(self, config: Config):
        self.cfg = config

    def construir(self, ruta_salida: str) -> str:
        """Construye la plantilla territorial de 28 slides."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # Slide 1: Portada
        self._slide_portada(prs, blank)
        # Slide 2: Portada área
        self._slide_portada_area(prs, blank)
        # Slide 3: Tendencias por eje
        self._slide_tendencias(prs, blank)
        # Slide 4: DOFA
        self._slide_dofa(prs, blank)
        # Slide 5: Divisor Hoja de Ruta
        self._slide_divisor(prs, blank)

        # Slides 6-24: Contribuciones a cada eje
        for eje_key in ORDEN_EJES:
            eje = EJES_REFERENCIA[eje_key]
            color = self.EJE_COLORS[eje_key]
            # Slide resumen del eje
            self._slide_resumen_eje(prs, blank, eje, color)
            # Slide por cada línea
            for le_idx in range(1, len(eje.lineas) + 1):
                self._slide_contribucion(prs, blank, eje, le_idx, color)

        # Slide 25: Metas
        self._slide_metas(prs, blank)
        # Slide 26: Presupuesto
        self._slide_presupuesto(prs, blank)
        # Slide 27: Gracias
        self._slide_gracias(prs, blank)

        prs.save(ruta_salida)
        log.info(f"✅ Plantilla territorial: {ruta_salida} "
                 f"({len(prs.slides)} slides)")
        return ruta_salida

    # --- Helpers ---

    def _bar(self, sl, y=0, h=0.7, c=None):
        c = c or self.C_AZ
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y),
                                Inches(13.333), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()

    def _txt(self, sl, l, t, w, h, tx, sz=14, b=False, c=None, al=PP_ALIGN.LEFT):
        c = c or self.C_GR
        tb = sl.shapes.add_textbox(Inches(l), Inches(t),
                                   Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tx
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = 'Calibri'
        p.alignment = al
        return tf

    def _lbl(self, sl, l, t, w, h, tx, bg=None):
        bg = bg or self.C_RJ
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        s.line.fill.background()
        tf = s.text_frame
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tx
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = self.C_BL
        r.font.name = 'Calibri'

    # --- Slides ---

    def _slide_portada(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = self.C_AZ
        self._txt(sl, 0.8, 1.8, 8, 1, 'SEGUIMIENTO {{TRIMESTRE}}',
                  44, True, self.C_BL)
        self._txt(sl, 0.8, 3.2, 10, 1, '{{UNIDAD}}',
                  28, True, self.C_BL)
        self._txt(sl, 0.8, 4.5, 8, 0.5,
                  'Planeación Estratégica {{AÑO}} — ProColombia',
                  16, False, self.C_GC)
        self._txt(sl, 0.8, 5.2, 4, 0.4, '{{TIPO}}',
                  14, False, self.C_GC)

    def _slide_portada_area(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        self._bar(sl)
        self._txt(sl, 0.5, 0.1, 10, 0.5,
                  '{{TIPO}} — {{UNIDAD}}', 20, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3)
        self._txt(sl, 0.5, 3.0, 12, 1, '{{UNIDAD}}',
                  36, True, self.C_AZ, PP_ALIGN.CENTER)

    def _slide_tendencias(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        self._bar(sl)
        self._txt(sl, 0.5, 0.1, 10, 0.5,
                  'TENDENCIAS DESDE {{UNIDAD}}', 18, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3)

        # 3 columnas: TUR, INV, EXP
        cols = [
            ('TUR', 'TURISMO', '{{TEND_TUR}}', '{{FOCO_TUR}}',
             '{{APORTE_TUR}}', self.EJE_COLORS['TUR']),
            ('INV', 'INVERSIÓN', '{{TEND_INV}}', '{{FOCO_INV}}',
             '{{APORTE_INV}}', self.EJE_COLORS['INV']),
            ('EXP', 'EXPORTACIONES', '{{TEND_EXP}}', '{{FOCO_EXP}}',
             '{{APORTE_EXP}}', self.EJE_COLORS['EXP']),
        ]
        for i, (key, nombre, t_tend, t_foco, t_aporte, col) in enumerate(cols):
            x = 0.3 + i * 4.3
            self._lbl(sl, x, 0.85, 4.0, 0.3, nombre, col)
            self._txt(sl, x, 1.25, 4.0, 0.2, 'Tendencias:', 8, True, self.C_AZ)
            self._txt(sl, x, 1.45, 4.0, 1.5, t_tend, 7, False, self.C_GR)
            self._txt(sl, x, 3.1, 4.0, 0.2, 'Foco / Prioridades:', 8, True, self.C_AZ)
            self._txt(sl, x, 3.3, 4.0, 1.5, t_foco, 7, False, self.C_GR)
            self._txt(sl, x, 5.0, 4.0, 0.2, 'Aporte:', 8, True, self.C_AZ)
            self._txt(sl, x, 5.2, 4.0, 1.5, t_aporte, 7, False, self.C_GR)

    def _slide_dofa(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        self._bar(sl)
        self._txt(sl, 0.5, 0.1, 10, 0.5,
                  'DOFA — {{UNIDAD}}', 18, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3)
        cpos = {
            'DEB': ('DEBILIDADES', 0.5, 1.2, 5.8, 2.5, self.C_RJ),
            'OPO': ('OPORTUNIDADES', 6.8, 1.2, 5.8, 2.5, self.C_AZ2),
            'FOR': ('FORTALEZAS', 0.5, 4.2, 5.8, 2.5, self.C_VD),
            'AME': ('AMENAZAS', 6.8, 4.2, 5.8, 2.5, self.C_NJ),
        }
        for pre, (nombre, l, t, w, h, c) in cpos.items():
            self._lbl(sl, l, t - 0.35, w, 0.3, nombre, c)
            items = '\n'.join(f'{{{{{pre}_BASE_{i}}}}}' for i in range(1, 11))
            self._txt(sl, l, t, w, h, items, 7, False, self.C_GR)

    def _slide_divisor(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = self.C_AZ
        self._txt(sl, 1, 2.5, 11, 1, 'HOJA DE RUTA', 44, True,
                  self.C_BL, PP_ALIGN.CENTER)
        self._txt(sl, 1, 3.8, 11, 0.5, '{{UNIDAD}}', 24, False,
                  self.C_GC, PP_ALIGN.CENTER)

    def _slide_resumen_eje(self, prs, blank, eje: EjeMisional,
                           color: RGBColor):
        sl = prs.slides.add_slide(blank)
        self._bar(sl, c=color)
        self._txt(sl, 0.5, 0.1, 10, 0.5,
                  f'{eje.nombre} — LÍNEAS ESTRATÉGICAS',
                  18, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3, c=color)

        for i, _ in enumerate(eje.lineas, 1):
            y = 1.0 + (i - 1) * 1.4
            tag = f'{{{{{eje.prefijo}_LE{i}_NOMBRE}}}}'
            self._txt(sl, 0.5, y, 0.6, 0.6, f'{i}.', 24, True, color)
            self._txt(sl, 1.3, y, 11, 0.8, tag, 11, False, self.C_GR)

    def _slide_contribucion(self, prs, blank, eje: EjeMisional,
                            le_idx: int, color: RGBColor):
        pre = eje.prefijo
        sl = prs.slides.add_slide(blank)
        self._bar(sl, c=color)
        self._txt(sl, 0.5, 0.1, 12, 0.5,
                  f'CONTRIBUCIÓN A {eje.nombre}', 14, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3, c=color)

        self._lbl(sl, 0.5, 0.85, 12.3, 0.3,
                  f'LÍNEA ESTRATÉGICA {le_idx}', color)
        self._txt(sl, 0.5, 1.25, 12.3, 0.5,
                  f'{{{{{pre}_LE{le_idx}_NOMBRE}}}}', 10, False, self.C_AZ2)

        # Acciones
        self._lbl(sl, 0.5, 1.85, 6, 0.3,
                  'ACCIONES {{TIPO}} — {{TRIMESTRE}}', self.C_AZ)
        self._lbl(sl, 6.8, 1.85, 6, 0.3,
                  'ACTIVIDADES CLAVE — {{TRIMESTRE}}', self.C_AZ)

        for i in range(1, 10):
            y = 2.25 + (i - 1) * 0.5
            if y > 5.0:
                break
            self._txt(sl, 0.5, y, 6.0, 0.25,
                      f'{{{{{pre}_LE{le_idx}_ACC_{i}}}}}',
                      8, True, RGBColor(0x33, 0x33, 0x33))
            self._txt(sl, 6.8, y, 6.0, 0.25,
                      f'{{{{{pre}_LE{le_idx}_ACT_{i}}}}}',
                      8, False, self.C_GR)
            self._txt(sl, 0.5, y + 0.2, 12.3, 0.25,
                      f'{{{{{pre}_LE{le_idx}_AVA_{i}}}}}',
                      7, False, self.C_GR)

        # Indicadores
        self._lbl(sl, 0.5, 5.5, 12.3, 0.3,
                  f'INDICADORES — {{{{AÑO}}}}', self.C_AZ)
        for i in range(1, 5):
            y = 5.9 + (i - 1) * 0.3
            self._txt(sl, 0.5, y, 4, 0.25,
                      f'{{{{{pre}_LE{le_idx}_IND_{i}}}}}', 7, False, self.C_GR)
            self._txt(sl, 4.8, y, 2, 0.25,
                      f'{{{{{pre}_LE{le_idx}_META_{i}}}}}', 7, False, self.C_GR)
            self._txt(sl, 7.0, y, 2, 0.25,
                      f'{{{{{pre}_LE{le_idx}_RES_{i}}}}}', 7, False, self.C_GR)

    def _slide_metas(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        self._bar(sl)
        self._txt(sl, 0.5, 0.1, 10, 0.5,
                  'METAS GENERALES — {{UNIDAD}}', 18, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3)

        rows = 11
        ts = sl.shapes.add_table(rows, 4, Inches(0.5), Inches(1.0),
                                 Inches(12.3), Inches(0.4 * rows))
        tb = ts.table
        for i, h in enumerate(['Indicador', 'Meta', 'Avance', '% Avance']):
            cl = tb.cell(0, i)
            cl.text = h
            for p in cl.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self.C_BL
            cl.fill.solid()
            cl.fill.fore_color.rgb = self.C_AZ
        for r in range(1, rows):
            tb.cell(r, 0).text = f'{{{{META_{r}_IND}}}}'
            tb.cell(r, 1).text = f'{{{{META_{r}_META}}}}'
            tb.cell(r, 2).text = f'{{{{META_{r}_AVA}}}}'
            tb.cell(r, 3).text = ''
            for c in range(4):
                for p in tb.cell(r, c).text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.name = 'Calibri'

    def _slide_presupuesto(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        self._bar(sl)
        self._txt(sl, 0.5, 0.1, 10, 0.5,
                  'PRESUPUESTO — {{UNIDAD}}', 18, True, self.C_BL)
        self._bar(sl, y=7.2, h=0.3)
        self._txt(sl, 0.5, 1.5, 12, 4, '{{PRESUPUESTO}}',
                  10, False, self.C_GR)

    def _slide_gracias(self, prs, blank):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = self.C_AZ
        self._txt(sl, 1, 2.5, 11, 1.5, 'GRACIAS', 54, True,
                  self.C_BL, PP_ALIGN.CENTER)
        self._txt(sl, 1, 4.2, 11, 0.5,
                  '{{UNIDAD}} — ProColombia {{AÑO}}',
                  18, False, self.C_GC, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# GENERADOR DE PPTX TERRITORIAL
# ═══════════════════════════════════════════════════════════════════════



# ═══════════════════════════════════════════════════════════════════════
# GENERADOR PPTX — TERRITORIAL
# ═══════════════════════════════════════════════════════════════════════

class GeneradorPPTXTerritorial:
    """
    Genera presentaciones PPTX para unidades TERRITORIALES
    usando la plantilla territorial con marcadores {{...}}.
    """

    def __init__(self, config: Config):
        self.cfg = config

    def _construir_reemplazos(self, data: Dict) -> Dict[str, str]:
        """Construye dict marcador→valor para territorial."""
        cfg_data = data['config']
        q = cfg_data.get('Trimestre en seguimiento', 'Q1')
        unidad = cfg_data.get('Nombre de la unidad', 'Área')
        año = cfg_data.get('Año', '2026')
        tipo = cfg_data.get('Tipo de unidad', 'HUB')

        r = {
            '{{TRIMESTRE}}': q, '{{UNIDAD}}': unidad,
            '{{AÑO}}': año, '{{TIPO}}': tipo,
            '{{PRESUPUESTO}}': data.get('presupuesto_texto', ''),
        }

        # DOFA (mismos marcadores que misional)
        prefix_map = {'DEBILIDADES': 'DEB', 'OPORTUNIDADES': 'OPO',
                      'FORTALEZAS': 'FOR', 'AMENAZAS': 'AME'}
        for cuadrante, prefix in prefix_map.items():
            items = data.get('dofa', {}).get(cuadrante, [])
            for i in range(1, 11):
                tag = f'{{{{{prefix}_BASE_{i}}}}}'
                if i <= len(items):
                    item = items[i - 1]
                    if item.get('estado') == 'Se elimina':
                        r[tag] = ''
                    elif item.get('estado') == 'Se actualiza' \
                            and item.get('actualizacion'):
                        r[tag] = _trunc(item['actualizacion'], 200)
                    else:
                        r[tag] = _trunc(item.get('base', ''), 200)
                else:
                    r[tag] = ''

        # Tendencias por eje
        tend_data = data.get('tendencias_por_eje', {})
        for eje_key in ['TUR', 'INV', 'EXP']:
            eje_tend = tend_data.get(eje_key, {})
            for sec_key, marker_key in [
                ('tendencias', 'TEND'), ('foco', 'FOCO'), ('aporte', 'APORTE')
            ]:
                items = eje_tend.get(sec_key, [])
                texto = '\n'.join(f'• {it}' for it in items) if items else ''
                r[f'{{{{{marker_key}_{eje_key}}}}}'] = _trunc(texto, 800)

        # Contribuciones a cada eje
        contribuciones = data.get('contribuciones', {})
        for eje_key in ORDEN_EJES:
            eje = EJES_REFERENCIA[eje_key]
            pre = eje.prefijo
            contrib_list = contribuciones.get(eje_key, [])

            for le_idx in range(1, len(eje.lineas) + 1):
                # Nombre de la línea (siempre desde referencia)
                r[f'{{{{{pre}_LE{le_idx}_NOMBRE}}}}'] = \
                    _trunc(eje.lineas[le_idx - 1], 300)

                le_data = contrib_list[le_idx - 1] \
                    if le_idx <= len(contrib_list) else None

                for acc_idx in range(1, self.cfg.max_acciones_por_linea + 1):
                    for suffix in ['ACC', 'ACT', 'AVA']:
                        tag = f'{{{{{pre}_LE{le_idx}_{suffix}_{acc_idx}}}}}'
                        if le_data and acc_idx <= len(le_data.get('acciones', [])):
                            acc = le_data['acciones'][acc_idx - 1]
                            val_map = {'ACC': 'accion', 'ACT': 'actividad',
                                       'AVA': 'avance'}
                            r[tag] = _trunc(acc.get(val_map[suffix], ''), 300)
                        else:
                            r[tag] = ''

                for ind_idx in range(1, self.cfg.max_indicadores_por_linea + 1):
                    for suffix, key in [('IND', 'indicador'),
                                        ('META', 'meta'), ('RES', 'avance')]:
                        tag = f'{{{{{pre}_LE{le_idx}_{suffix}_{ind_idx}}}}}'
                        if le_data and ind_idx <= len(
                                le_data.get('indicadores', [])):
                            r[tag] = str(le_data['indicadores'][
                                ind_idx - 1].get(key, ''))
                        else:
                            r[tag] = ''

        # Metas
        metas = data.get('metas', [])
        for i in range(1, self.cfg.max_metas + 1):
            for suffix, key in [('IND', 'indicador'),
                                ('META', 'meta'), ('AVA', 'avance')]:
                tag = f'{{{{META_{i}_{suffix}}}}}'
                if i <= len(metas):
                    r[tag] = str(metas[i - 1].get(key, ''))
                else:
                    r[tag] = ''
        return r

    def generar(self, data: Dict, ruta_plantilla: str,
                ruta_salida: str) -> str:
        """Genera PPTX territorial reemplazando marcadores."""
        log.info(f"📊 Plantilla territorial: {os.path.basename(ruta_plantilla)}")
        prs = Presentation(ruta_plantilla)
        reemplazos = self._construir_reemplazos(data)
        log.info(f"  Marcadores definidos: {len(reemplazos)}")

        gen = GeneradorPPTXMisional(self.cfg)
        total = 0
        for sl in prs.slides:
            total += gen._reemplazar_en_slide(sl, reemplazos)
        log.info(f"  Marcadores reemplazados: {total}")

        prs.save(ruta_salida)
        log.info(f"  ✅ PPTX territorial: {os.path.basename(ruta_salida)} "
                 f"({len(prs.slides)} slides)")
        return ruta_salida



# ═══════════════════════════════════════════════════════════════════════
# CONSTRUCTOR PLANTILLA PPTX — TRANSVERSAL
# ═══════════════════════════════════════════════════════════════════════

class ConstructorPlantillaTransversal:
    """
    Construye la plantilla PPTX transversal combinando:
    - Bloque propio (portada, DOFA, resumen LE propias, slides LE propias)
    - Bloque contribuciones (resumen + slides por cada eje)
    - Metas, presupuesto, cierre
    """

    def __init__(self, config: Config):
        self.cfg = config
        self._terr = ConstructorPlantillaTerritorial(config)

    def construir(self, ruta_salida: str) -> str:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        T = self._terr  # shortcut

        # 1: Portada
        T._slide_portada(prs, blank)
        # 2: Portada área
        T._slide_portada_area(prs, blank)
        # 3: DOFA
        T._slide_dofa(prs, blank)
        # 4: Resumen líneas propias
        self._slide_resumen_propias(prs, blank)
        # 5-9: Líneas propias LE1-LE5
        for i in range(1, self.cfg.max_lineas_estrategicas + 1):
            self._slide_linea_propia(prs, blank, i)
        # 10: Divisor contribuciones
        self._slide_divisor_contrib(prs, blank)
        # 11-29: Contribuciones (resumen + slides por eje)
        for eje_key in ORDEN_EJES:
            eje = EJES_REFERENCIA[eje_key]
            color = T.EJE_COLORS[eje_key]
            T._slide_resumen_eje(prs, blank, eje, color)
            for le_idx in range(1, len(eje.lineas) + 1):
                T._slide_contribucion(prs, blank, eje, le_idx, color)
        # 30: Metas
        T._slide_metas(prs, blank)
        # 31: Presupuesto
        T._slide_presupuesto(prs, blank)
        # 32: Gracias
        T._slide_gracias(prs, blank)

        prs.save(ruta_salida)
        log.info(f"✅ Plantilla transversal: {ruta_salida} "
                 f"({len(prs.slides)} slides)")
        return ruta_salida

    def _slide_resumen_propias(self, prs, blank):
        T = self._terr
        sl = prs.slides.add_slide(blank)
        T._bar(sl)
        T._txt(sl, 0.5, 0.1, 10, 0.5,
               '{{UNIDAD}} — LÍNEAS ESTRATÉGICAS PROPIAS',
               18, True, T.C_BL)
        T._bar(sl, y=7.2, h=0.3)
        for i in range(1, self.cfg.max_lineas_estrategicas + 1):
            y = 1.0 + (i - 1) * 1.1
            T._txt(sl, 0.5, y, 0.6, 0.6, f'{i}.', 24, True, T.C_AZ)
            T._txt(sl, 1.3, y, 11, 0.8,
                   f'{{{{LE{i}_NOMBRE}}}}', 11, False, T.C_GR)

    def _slide_linea_propia(self, prs, blank, le_idx: int):
        T = self._terr
        sl = prs.slides.add_slide(blank)
        T._bar(sl)
        T._txt(sl, 0.5, 0.1, 12, 0.5,
               'HOJA DE RUTA ({{UNIDAD}})', 14, True, T.C_BL)
        T._bar(sl, y=7.2, h=0.3)
        T._lbl(sl, 0.5, 0.85, 12.3, 0.3,
               f'LÍNEA ESTRATÉGICA {le_idx}', T.C_RJ)
        T._txt(sl, 0.5, 1.25, 12.3, 0.5,
               f'{{{{LE{le_idx}_NOMBRE}}}}', 10, False, T.C_AZ2)
        T._lbl(sl, 0.5, 1.85, 6, 0.3,
               'ACCIONES — {{TRIMESTRE}}', T.C_AZ)
        T._lbl(sl, 6.8, 1.85, 6, 0.3,
               'ACTIVIDADES CLAVE', T.C_AZ)
        for i in range(1, 10):
            y = 2.25 + (i - 1) * 0.5
            if y > 5.0:
                break
            T._txt(sl, 0.5, y, 6.0, 0.25,
                   f'{{{{LE{le_idx}_ACC_{i}}}}}', 8, True,
                   RGBColor(0x33, 0x33, 0x33))
            T._txt(sl, 6.8, y, 6.0, 0.25,
                   f'{{{{LE{le_idx}_ACT_{i}}}}}', 8, False, T.C_GR)
            T._txt(sl, 0.5, y + 0.2, 12.3, 0.25,
                   f'{{{{LE{le_idx}_AVA_{i}}}}}', 7, False, T.C_GR)
        T._lbl(sl, 0.5, 5.5, 12.3, 0.3,
               f'INDICADORES — {{{{AÑO}}}}', T.C_AZ)
        for i in range(1, 5):
            y = 5.9 + (i - 1) * 0.3
            T._txt(sl, 0.5, y, 4, 0.25,
                   f'{{{{LE{le_idx}_IND_{i}}}}}', 7, False, T.C_GR)
            T._txt(sl, 4.8, y, 2, 0.25,
                   f'{{{{LE{le_idx}_META_{i}}}}}', 7, False, T.C_GR)
            T._txt(sl, 7.0, y, 2, 0.25,
                   f'{{{{LE{le_idx}_RES_{i}}}}}', 7, False, T.C_GR)

    def _slide_divisor_contrib(self, prs, blank):
        T = self._terr
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = T.C_AZ
        T._txt(sl, 1, 2.0, 11, 1, 'CONTRIBUCIONES A', 36, True,
               T.C_BL, PP_ALIGN.CENTER)
        T._txt(sl, 1, 3.0, 11, 1, 'EJES MISIONALES', 44, True,
               T.C_BL, PP_ALIGN.CENTER)
        T._txt(sl, 1, 4.5, 11, 0.5, '{{UNIDAD}}', 20, False,
               T.C_GC, PP_ALIGN.CENTER)



# ═══════════════════════════════════════════════════════════════════════
# GENERADOR PPTX — TRANSVERSAL
# ═══════════════════════════════════════════════════════════════════════

class GeneradorPPTXTransversal:
    """
    Genera PPTX transversal combinando reemplazos de:
    - Líneas propias (misma lógica que misional: LE1..LE5)
    - Contribuciones a ejes (misma lógica que territorial: MP_LE1..EXP_LE4)
    """

    def __init__(self, config: Config):
        self.cfg = config
        self._gen_misional = GeneradorPPTXMisional(config)
        self._gen_territorial = GeneradorPPTXTerritorial(config)

    def generar(self, data: Dict, ruta_plantilla: str,
                ruta_salida: str) -> str:
        log.info(f"📊 Plantilla transversal: "
                 f"{os.path.basename(ruta_plantilla)}")
        prs = Presentation(ruta_plantilla)

        # Combinar reemplazos de ambas fuentes
        r_misional = self._gen_misional._construir_reemplazos(data)
        r_territorial = self._gen_territorial._construir_reemplazos(data)

        # Territorial sobreescribe globals (TRIMESTRE, UNIDAD, etc.)
        # pero misional tiene LE1..LE5 propias que territorial no tiene
        reemplazos = {**r_territorial, **r_misional}
        # Restaurar contribuciones territoriales (MP_*, TUR_*, etc.)
        for k, v in r_territorial.items():
            if any(k.startswith(f'{{{{{p}_') for p in ['MP', 'TUR', 'INV', 'EXP']):
                reemplazos[k] = v

        log.info(f"  Marcadores combinados: {len(reemplazos)}")

        # Determinar líneas propias con contenido para eliminación dinámica
        lineas_con = set()
        for i in range(1, self.cfg.max_lineas_estrategicas + 1):
            if i <= len(data.get('lineas', [])):
                if self._gen_misional._linea_tiene_contenido(
                        data['lineas'][i - 1]):
                    lineas_con.add(i)

        slides_a_eliminar = []
        for si, slide in enumerate(prs.slides):
            le_en = set()
            for n in range(1, self.cfg.max_lineas_estrategicas + 1):
                if self._gen_misional._slide_tiene_marcadores_le(slide, n):
                    le_en.add(n)
            # Solo eliminar si tiene marcadores LE propios y ninguno activo
            # No eliminar contribuciones (MP_LE, TUR_LE, etc.)
            if le_en and le_en.isdisjoint(lineas_con):
                # Verificar que no sean contribuciones
                full = ''.join(
                    ''.join(r.text for r in p.runs)
                    for shape in self._gen_misional._iter_shapes(slide.shapes)
                    if shape.has_text_frame
                    for p in shape.text_frame.paragraphs
                )
                has_contrib = any(f'{p}_LE' in full
                                 for p in ['MP', 'TUR', 'INV', 'EXP'])
                if not has_contrib:
                    slides_a_eliminar.append(si)

        # Reemplazar
        total = sum(
            self._gen_misional._reemplazar_en_slide(sl, reemplazos)
            for sl in prs.slides
        )
        log.info(f"  Marcadores reemplazados: {total}")

        if slides_a_eliminar:
            log.info(f"  Eliminando {len(slides_a_eliminar)} slides "
                     f"de líneas propias vacías")
            self._gen_misional._eliminar_slides(prs, slides_a_eliminar)

        prs.save(ruta_salida)
        log.info(f"  ✅ PPTX transversal: {os.path.basename(ruta_salida)} "
                 f"({len(prs.slides)} slides)")
        return ruta_salida
